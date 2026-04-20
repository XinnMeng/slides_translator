from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from pptx_translator.extractor import extract_text_items
from pptx_translator.pipeline import run_translation
from pptx_translator.translator import TranslatorBackend


class FakeTranslator(TranslatorBackend):
    def translate_text(self, text: str, source_lang: str, target_lang: str) -> str:
        replacement = {
            "Hallo": "Hello",
            "Welt": "World",
            "Guten Morgen": "Good morning",
        }
        out = text
        for src, dst in replacement.items():
            out = out.replace(src, dst)
        return out


def test_tiny_end_to_end_multiple_slides(tmp_path: Path) -> None:
    in_file = tmp_path / "in.pptx"
    out_file = tmp_path / "out.pptx"

    prs = Presentation()

    slide_1 = prs.slides.add_slide(prs.slide_layouts[5])
    box_1 = slide_1.shapes.add_textbox(left=0, top=0, width=2_000_000, height=1_000_000)
    box_1.text_frame.paragraphs[0].text = "Hallo Welt"

    slide_2 = prs.slides.add_slide(prs.slide_layouts[5])
    box_2 = slide_2.shapes.add_textbox(left=0, top=0, width=2_000_000, height=1_000_000)
    box_2.text_frame.paragraphs[0].text = "Guten Morgen"

    prs.save(in_file)

    result = run_translation(
        input_path=in_file,
        output_path=out_file,
        source_lang="de",
        target_lang="en",
        translator=FakeTranslator(),
        dry_run=False,
    )

    assert result["mode"] == "translate"
    assert result["applied"] == 2

    out_prs = Presentation(str(out_file))
    texts = [item.text for item in extract_text_items(out_prs)]
    assert "Hello World" in texts
    assert "Good morning" in texts
