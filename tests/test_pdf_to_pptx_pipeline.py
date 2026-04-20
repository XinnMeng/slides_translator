from __future__ import annotations

from pathlib import Path

import fitz
from pptx import Presentation

from pptx_translator.pipelines.pdf_to_pptx_pipeline import run_pdf_to_pptx_translation
from pptx_translator.translator import TranslatorBackend


class FakeJaToZhTranslator(TranslatorBackend):
    def translate_text(self, text: str, source_lang: str, target_lang: str) -> str:
        return f"中文:{text}"


def _make_multipage_pdf(path: Path) -> None:
    doc = fitz.open()
    p1 = doc.new_page(width=500, height=400)
    p1.insert_text((50, 80), "日本語テキスト1")

    p2 = doc.new_page(width=500, height=400)
    p2.insert_text((50, 80), "日本語テキスト2")

    doc.save(str(path))
    doc.close()


def test_pdf_to_pptx_multi_page(tmp_path: Path) -> None:
    pdf_path = tmp_path / "in.pdf"
    out_path = tmp_path / "out.pptx"
    _make_multipage_pdf(pdf_path)

    result = run_pdf_to_pptx_translation(
        input_pdf=pdf_path,
        output_pptx=out_path,
        translator=FakeJaToZhTranslator(),
        source_lang="ja",
        target_lang="zh-CN",
    )

    assert result["pages"] == 2

    prs = Presentation(str(out_path))
    assert len(prs.slides) == 2

    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                all_text.append(shape.text_frame.text)

    assert any("中文:" in t for t in all_text)
