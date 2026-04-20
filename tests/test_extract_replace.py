from __future__ import annotations

from pptx import Presentation

from pptx_translator.extractor import extract_text_items
from pptx_translator.models import text_item_id
from pptx_translator.replacer import apply_translations


def make_sample_presentation() -> Presentation:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    box = slide.shapes.add_textbox(left=0, top=0, width=3_000_000, height=1_000_000)
    tf = box.text_frame
    tf.paragraphs[0].text = "Hallo Welt"
    p2 = tf.add_paragraph()
    p2.text = "Already English"

    table = slide.shapes.add_table(rows=1, cols=1, left=0, top=1_200_000, width=3_000_000, height=1_000_000).table
    table.cell(0, 0).text = "Guten Morgen"

    return prs


def test_extract_text_items_from_shapes_and_table() -> None:
    prs = make_sample_presentation()

    items = extract_text_items(prs)
    texts = [i.text for i in items]

    assert "Hallo Welt" in texts
    assert "Already English" in texts
    assert "Guten Morgen" in texts
    assert len(items) == 3


def test_replace_translated_text_items() -> None:
    prs = make_sample_presentation()
    items = extract_text_items(prs)

    translated = {}
    for item in items:
        key = text_item_id(item.ref)
        if item.text == "Hallo Welt":
            translated[key] = "Hello world"
        elif item.text == "Guten Morgen":
            translated[key] = "Good morning"
        else:
            translated[key] = item.text

    applied = apply_translations(prs, items, translated)
    assert applied == 2

    new_items = extract_text_items(prs)
    new_texts = [i.text for i in new_items]
    assert "Hello world" in new_texts
    assert "Good morning" in new_texts
    assert "Already English" in new_texts


def test_extract_text_items_multiple_slides() -> None:
    prs = Presentation()

    slide_1 = prs.slides.add_slide(prs.slide_layouts[5])
    box_1 = slide_1.shapes.add_textbox(left=0, top=0, width=2_000_000, height=1_000_000)
    box_1.text_frame.paragraphs[0].text = "Hallo"

    slide_2 = prs.slides.add_slide(prs.slide_layouts[5])
    box_2 = slide_2.shapes.add_textbox(left=0, top=0, width=2_000_000, height=1_000_000)
    box_2.text_frame.paragraphs[0].text = "Guten Morgen"

    items = extract_text_items(prs)
    assert len(items) == 2

    by_slide = {}
    for item in items:
        by_slide[item.ref.slide_index] = by_slide.get(item.ref.slide_index, 0) + 1

    assert by_slide == {0: 1, 1: 1}
