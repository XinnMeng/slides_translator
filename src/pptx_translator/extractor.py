from __future__ import annotations

import logging
from collections.abc import Iterable

from pptx.presentation import Presentation

from .models import TextItem, TextItemRef

logger = logging.getLogger(__name__)


def _paragraph_text(paragraph) -> str:
    text = "".join(run.text for run in paragraph.runs) if paragraph.runs else paragraph.text
    return text.strip()


def extract_text_items(prs: Presentation) -> list[TextItem]:
    """Extract candidate text snippets from text frames and table cells."""
    items: list[TextItem] = []

    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            # Text boxes/placeholders/other text-frame-capable shapes.
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                text_frame = shape.text_frame
                for p_idx, paragraph in enumerate(text_frame.paragraphs):
                    text = _paragraph_text(paragraph)
                    if not text:
                        continue
                    items.append(
                        TextItem(
                            ref=TextItemRef(
                                slide_index=slide_idx,
                                shape_index=shape_idx,
                                text_frame_paragraph_index=p_idx,
                            ),
                            text=text,
                        )
                    )

            # Tables
            if getattr(shape, "has_table", False) and shape.has_table:
                table = shape.table
                for r_idx, row in enumerate(table.rows):
                    for c_idx, cell in enumerate(row.cells):
                        for p_idx, paragraph in enumerate(cell.text_frame.paragraphs):
                            text = _paragraph_text(paragraph)
                            if not text:
                                continue
                            items.append(
                                TextItem(
                                    ref=TextItemRef(
                                        slide_index=slide_idx,
                                        shape_index=shape_idx,
                                        table_row_index=r_idx,
                                        table_col_index=c_idx,
                                        table_paragraph_index=p_idx,
                                    ),
                                    text=text,
                                )
                            )

    logger.info("Extracted %s text items from %s slides", len(items), len(prs.slides))
    return items


def group_items_by_slide(items: Iterable[TextItem]) -> dict[int, list[TextItem]]:
    grouped: dict[int, list[TextItem]] = {}
    for item in items:
        grouped.setdefault(item.ref.slide_index, []).append(item)
    return grouped
