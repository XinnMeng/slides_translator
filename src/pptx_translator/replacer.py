from __future__ import annotations

import logging

from pptx.presentation import Presentation

from .extractor import group_items_by_slide
from .models import TextItem, text_item_id

logger = logging.getLogger(__name__)


def _replace_paragraph_text(paragraph, new_text: str, preferred_run_index: int | None) -> None:
    """
    Replace paragraph text while preserving paragraph properties and preferred-run style.

    Strategy:
    - If runs exist, write new text into preferred run (or first run) and clear other runs.
    - If there are no runs, set paragraph.text directly.
    """
    if paragraph.runs:
        target_index = preferred_run_index if preferred_run_index is not None else 0
        target_index = min(max(target_index, 0), len(paragraph.runs) - 1)

        paragraph.runs[target_index].text = new_text
        for idx, run in enumerate(paragraph.runs):
            if idx != target_index:
                run.text = ""
    else:
        paragraph.text = new_text


def apply_translations(prs: Presentation, items: list[TextItem], translated_by_id: dict[str, str]) -> int:
    """Apply translated text back into the presentation using extracted item references."""
    applied = 0

    for slide_idx, slide_items in sorted(group_items_by_slide(items).items()):
        logger.info("Applying %s items to slide %s", len(slide_items), slide_idx)

        for item in slide_items:
            item_key = text_item_id(item.ref)
            if item_key not in translated_by_id:
                logger.info("Skipped item ID %s because no translated result was returned", item_key)
                continue

            translated_text = translated_by_id[item_key]
            if translated_text == item.text:
                logger.info("Skipped item ID %s because translated text is unchanged", item_key)
                continue

            try:
                slide = prs.slides[item.ref.slide_index]
                shape = slide.shapes[item.ref.shape_index]

                if item.ref.text_frame_paragraph_index is not None:
                    paragraph = shape.text_frame.paragraphs[item.ref.text_frame_paragraph_index]
                    _replace_paragraph_text(paragraph, translated_text, item.ref.text_frame_run_index)
                    applied += 1
                    logger.info("Updated item ID %s", item_key)
                    continue

                if (
                    item.ref.table_row_index is not None
                    and item.ref.table_col_index is not None
                    and item.ref.table_paragraph_index is not None
                ):
                    cell = shape.table.rows[item.ref.table_row_index].cells[item.ref.table_col_index]
                    paragraph = cell.text_frame.paragraphs[item.ref.table_paragraph_index]
                    _replace_paragraph_text(paragraph, translated_text, item.ref.table_run_index)
                    applied += 1
                    logger.info("Updated item ID %s", item_key)
                    continue

                logger.info("Skipped item ID %s because reference type is unsupported", item_key)
            except Exception as exc:
                logger.warning("Skipped item ID %s because replacement failed: %s", item_key, exc)

    logger.info("Applied %s translated text items", applied)
    return applied
