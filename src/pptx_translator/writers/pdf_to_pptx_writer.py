from __future__ import annotations

import logging

from pptx import Presentation
from pptx.util import Emu

from ..extractors.pdf_extractor import PdfPageBlocks, PdfTextBlock

logger = logging.getLogger(__name__)


class PdfCoordinateMapper:
    def __init__(self, pdf_width: float, pdf_height: float, slide_width_emu: int, slide_height_emu: int) -> None:
        self.scale_x = slide_width_emu / max(pdf_width, 1.0)
        self.scale_y = slide_height_emu / max(pdf_height, 1.0)

    def map_bbox(self, bbox: tuple[float, float, float, float]) -> tuple[int, int, int, int]:
        x0, y0, x1, y1 = bbox
        left = int(x0 * self.scale_x)
        top = int(y0 * self.scale_y)
        width = max(int((x1 - x0) * self.scale_x), 1)
        height = max(int((y1 - y0) * self.scale_y), 1)
        return left, top, width, height


def write_pdf_pages_to_pptx(
    pages: list[PdfPageBlocks],
    translated_text_by_key: dict[tuple[int, int], str],
) -> Presentation:
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    for page in pages:
        slide = prs.slides.add_slide(blank_layout)
        mapper = PdfCoordinateMapper(
            pdf_width=page.page_width,
            pdf_height=page.page_height,
            slide_width_emu=int(prs.slide_width),
            slide_height_emu=int(prs.slide_height),
        )

        inserted = 0
        for block in page.blocks:
            key = (block.page_index, block.block_index)
            text = translated_text_by_key.get(key, block.text)
            left, top, width, height = mapper.map_bbox(block.bbox)

            textbox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
            textbox.text_frame.paragraphs[0].text = text
            inserted += 1

        logger.info("Text boxes inserted on slide %s: %s", page.page_index, inserted)

    return prs
